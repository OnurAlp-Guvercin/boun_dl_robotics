"""Build the VLM-Guided Robot Reaching presentation (.pptx)."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"

# ---- palette -----------------------------------------------------------------
NAVY = RGBColor(0x0F, 0x2C, 0x4A)
BLUE = RGBColor(0x2A, 0x7D, 0xE1)
ORANGE = RGBColor(0xE8, 0x59, 0x0C)
GREEN = RGBColor(0x16, 0x65, 0x34)
RUST = RGBColor(0x9A, 0x34, 0x12)
DARK = RGBColor(0x22, 0x2A, 0x33)
GREY = RGBColor(0x5B, 0x66, 0x70)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
PALE_O = RGBColor(0xFD, 0xEE, 0xE3)
PALE_G = RGBColor(0xE7, 0xF3, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SKY = RGBColor(0xBF, 0xD6, 0xEE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    shp = s.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, *rest) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.name = "Calibri"
            if rest and rest[0]:
                r.font.italic = True
    return tb


def header(s, title, kicker=None):
    rect(s, 0, 0, SW, Inches(1.1), NAVY)
    rect(s, 0, Inches(1.1), SW, Pt(4), ORANGE)
    txt(s, Inches(0.55), Inches(0.14), Inches(12.2), Inches(0.82),
        [[(title, 29, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        txt(s, Inches(0.58), Inches(0.0), Inches(12), Inches(0.30),
            [[(kicker, 11.5, SKY, True)]], anchor=MSO_ANCHOR.BOTTOM)


def bullets(s, items, x=Inches(0.7), y=Inches(1.5), w=Inches(12), h=Inches(5.6),
            size=18, gap=10):
    paras = []
    for level, text, *style in items:
        st = style[0] if style else ""
        bold = "b" in st
        color = ORANGE if "o" in st else (GREEN if "g" in st else (RUST if "r" in st else DARK))
        indent = "" if level <= 1 else "      "
        bullet = "▸  " if level == 0 else ("•  " if level == 1 else "–  ")
        paras.append([(indent + bullet + text, size - (level * 0.5), color, bold)])
    txt(s, x, y, w, h, paras, space_after=gap, line_spacing=1.05)


def footer(s, n):
    txt(s, Inches(0.4), Inches(7.06), Inches(8), Inches(0.33),
        [[("VLM-Guided Robot Reaching", 10, GREY, False)]])
    txt(s, Inches(12.2), Inches(7.06), Inches(0.9), Inches(0.33),
        [[(str(n), 10, GREY, True)]], align=PP_ALIGN.RIGHT)


def pic_fit(s, path, x, y, w, h):
    iw, ih = Image.open(path).size
    box_ar = w / h
    img_ar = iw / ih
    if img_ar > box_ar:
        nw, nh = w, int(w / img_ar)
    else:
        nh, nw = h, int(h * img_ar)
    return s.shapes.add_picture(str(path), x + (w - nw) // 2, y + (h - nh) // 2, nw, nh)


def panel_title(s, x, y, w, text, color):
    b = rect(s, x, y, w, Inches(0.5), color)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(17); r.font.bold = True
    r.font.color.rgb = WHITE


N = 0


def num():
    global N
    N += 1
    return N


# =============================================================================
# 1 — Title
# =============================================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.5), SW, Pt(4), ORANGE)
txt(s, Inches(0.9), Inches(1.55), Inches(11.6), Inches(2.1),
    [[("VLM-Guided Robot Reaching", 46, WHITE, True)],
     [("Language-conditioned tabletop manipulation with a "
       "vision-language model in the loop", 22, SKY, False)]],
    space_after=14)
txt(s, Inches(0.92), Inches(4.7), Inches(11.6), Inches(1.7),
    [[("A frozen Qwen3-VL detector feeds a behaviour-cloned navigation policy "
       "that drives a UR5e arm to touch a named object.", 17, WHITE, False)],
     [("Onuralp Güvercin   ·   Deep Learning in Robotics   ·   Final Project",
       15, SKY, False)]],
    space_after=10)

# =============================================================================
# 2 — Objective
# =============================================================================
s = slide()
header(s, "Objective", "WHAT ARE WE SOLVING?")
bullets(s, [
    (0, "Task: a robot arm must reach out and touch a specific object named in "
        "natural language — e.g. \"the red cube\".", "b"),
    (1, "Scene: a tabletop with 2–4 differently coloured boxes and spheres at "
        "random positions (6 colours × 2 shapes)."),
    (1, "The robot sees only a top-down camera image; it is never told the "
        "target's coordinates."),
    (0, "The central question", "ob"),
    (1, "Can a vision-language model act as the \"eyes\" that localise the "
        "target, while a small learned policy acts as the \"hands\" that move "
        "toward it?"),
    (0, "Success metric", "ob"),
    (1, "An episode succeeds if the end-effector makes contact with — or comes "
        "within 5 cm of — the correct target."),
    (1, "We report success rate, final EE-to-object distance, and steps taken."),
], y=Inches(1.5), gap=12, size=19)
footer(s, num())

# =============================================================================
# 3 — Architecture (diagram)
# =============================================================================
s = slide()
header(s, "System Architecture", "PERCEPTION → CONTROL → ACTION, IN A CLOSED LOOP")

box_y = Inches(1.95); box_h = Inches(1.45); bw = Inches(2.7); gap = Inches(0.45)
x0 = Inches(0.7)


def flowbox(x, title, sub, color):
    b = rect(s, x, box_y, bw, box_h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(11); r2.font.color.rgb = RGBColor(0xEA, 0xF1, 0xFB)


def arrow(x):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, box_y + Inches(0.52), gap, Inches(0.4))
    a.fill.solid(); a.fill.fore_color.rgb = GREY; a.line.fill.background(); a.shadow.inherit = False


pos = [x0 + i * (bw + gap) for i in range(4)]
flowbox(pos[0], "MuJoCo Env", "top-down image\n+ ee_pos", NAVY)
arrow(pos[0] + bw)
flowbox(pos[1], "Qwen3-VL (frozen)", "image + object name\n→ bbox (cx,cy,w,h)", BLUE)
arrow(pos[1] + bw)
flowbox(pos[2], "NavigationMLP", "[bbox, ee_pos]\n→ H × Δee (metres)", ORANGE)
arrow(pos[2] + bw)
flowbox(pos[3], "UR5e Arm", "apply H deltas\n→ contact?", NAVY)

txt(s, Inches(0.7), Inches(3.55), Inches(12), Inches(0.45),
    [[("Closed loop: the policy executes H steps, then re-queries the VLM   "
       "(H steps = 1 VLM call).", 15, GREY, False, True)]], align=PP_ALIGN.CENTER)

bullets(s, [
    (1, "Perception (VLM): localises the named object as a normalised bounding box. Used off-the-shelf, weights frozen."),
    (1, "Control (MLP): 7-D input  [bbox(4) + ee_pos(3)]  →  3·H end-effector deltas in metres."),
    (1, "The policy is trained purely by behaviour cloning on a scripted expert — no reinforcement learning, no VLM training."),
], y=Inches(4.15), gap=11, size=16.5)
footer(s, num())

# =============================================================================
# 4 — Method: data collection
# =============================================================================
s = slide()
header(s, "Method · 1 — Data Collection", "TEACHING WITH A SCRIPTED EXPERT")
bullets(s, [
    (0, "An automated expert generates the demonstrations — no human teleop.", "b"),
    (1, "Each scene spawns 2–4 objects via rejection sampling (≥ 10 cm apart, unique colours)."),
    (1, "For every object, a straight-line controller drives the end-effector "
        "toward it in small 4 cm steps."),
    (0, "What is recorded at each step", "ob"),
    (1, "the top-down RGB image (128×128),"),
    (1, "the end-effector position ee_pos (xyz),"),
    (1, "the ground-truth bounding box, computed by projecting the object "
        "through the camera matrix."),
    (0, "Only successful reaches are saved", "ob"),
    (1, "contact is detected from MuJoCo geom collisions between gripper and target;"),
    (1, "collection is parallelised across worker processes, each with its own seed."),
], y=Inches(1.45), gap=8, size=17)
footer(s, num())

# =============================================================================
# 5 — Method: policy + behaviour cloning
# =============================================================================
s = slide()
header(s, "Method · 2 — Policy & Behaviour Cloning", "A SMALL RESIDUAL MLP THAT OUTPUTS MOTION")
bullets(s, [
    (0, "NavigationMLP — a residual MLP (width 512, 4 residual blocks, SiLU + LayerNorm).", "b"),
    (1, "Input  (7-D):  bbox(cx,cy,w,h)  +  normalised ee_pos(x,y,z)."),
    (1, "Output  (3·H):  end-effector deltas Δee in metres, H steps ahead."),
    (0, "Trained by behaviour cloning (supervised regression)", "ob"),
    (1, "Target = the expert's actual step-to-step Δee; loss = MSE on the deltas."),
    (1, "AdamW + warmup→cosine LR; the output head is zero-initialised so the "
        "policy starts from gentle, near-zero motions."),
    (0, "Honest evaluation", "ob"),
    (1, "Stratified train/val/test split by colour+shape; samples from one "
        "trajectory never leak across splits."),
    (1, "Test RMSE ≈ 0.0007 m — the policy reproduces expert deltas almost exactly."),
], y=Inches(1.45), gap=8, size=17)
footer(s, num())

# =============================================================================
# 6 — Critical implementation details
# =============================================================================
s = slide()
header(s, "Method · 3 — Critical Design Choices", "THE DECISIONS THAT MAKE IT WORK")
bullets(s, [
    (0, "Fixed bounding box per episode.", "b"),
    (1, "The bbox is read once and held constant; the loop closes on ee_pos, not on "
        "a re-detected box. This makes behaviour robust to frame-to-frame VLM jitter."),
    (0, "Delta actions with clamping (±5 cm/step).", "b"),
    (1, "Predicting relative motion (not absolute targets) keeps every step bounded "
        "and stable — no large jumps."),
    (0, "Configurable horizon H (1–5).", "b"),
    (1, "H = how many steps run per VLM query. Small H = more re-planning (costly, "
        "accurate); large H = more open-loop drift between queries."),
    (0, "Robust VLM parsing + fallback.", "b"),
    (1, "Strips <think> reasoning, auto-detects 0–1000 / 0–128 / 0–1 coordinate "
        "scales, and falls back to a GT bbox if a call fails."),
], y=Inches(1.45), gap=9, size=17)
footer(s, num())

# =============================================================================
# 7 — Why VLM not trained + GRPO motivation
# =============================================================================
s = slide()
header(s, "We Did Not Train the VLM", "A SCOPE DECISION — AND THE PATH WE WOULD TAKE")
bullets(s, [
    (0, "In this project Qwen3-VL-4B is used frozen, off-the-shelf, as the detector.", "b"),
    (1, "Fine-tuning a 4-billion-parameter vision-language model needs multi-GPU "
        "memory and long training runs — compute we did not have."),
    (1, "So our contribution is the end-to-end pipeline + the control policy, with "
        "the VLM as a plug-in perception module."),
    (0, "If compute had been available", "ob"),
    (1, "We would fine-tune the VLM with GRPO (Group Relative Policy Optimization)."),
    (1, "Reward = how well the predicted bbox matches the target "
        "(IoU / centre error), or even downstream reaching success."),
    (1, "This adapts the detector to our exact camera, lighting, and object set — "
        "directly attacking the VLM-vs-GT gap shown in the results."),
], y=Inches(1.5), gap=12, size=18)
footer(s, num())

# =============================================================================
# 8 — GRPO diagram (PPO vs GRPO visual)
# =============================================================================
s = slide()
header(s, "PPO vs GRPO — The Idea", "GRPO DROPS THE CRITIC AND USES A GROUP BASELINE")
pic_fit(s, ASSETS / "grpo_diagram.png", Inches(0.4), Inches(1.35),
        Inches(8.6), Inches(5.4))
txt(s, Inches(0.5), Inches(6.75), Inches(8.4), Inches(0.4),
    [[("PPO trains an extra Value Model (\"critic\"); GRPO replaces it by sampling a "
       "group and normalising rewards.   (DeepSeekMath, 2024)", 10.5, GREY, False, True)]])
bullets(s, [
    (0, "PPO (top)", "ob"),
    (1, "needs a separate Value "
        "Model to estimate the "
        "baseline (advantage via GAE)."),
    (1, "= a second large network "
        "to train and store."),
    (0, "GRPO (bottom)", "ob"),
    (1, "samples a GROUP of G "
        "answers per prompt,"),
    (1, "scores each with the "
        "reward model,"),
    (1, "uses the group's mean/std "
        "as the baseline — no "
        "critic at all."),
], x=Inches(9.15), y=Inches(1.5), w=Inches(3.9), gap=8, size=14)
footer(s, num())

# =============================================================================
# 9 — PPO vs GRPO loss functions
# =============================================================================
s = slide()
header(s, "PPO vs GRPO — The Loss Functions", "SAME CLIPPED OBJECTIVE, DIFFERENT ADVANTAGE")

# shared ratio
txt(s, Inches(0.55), Inches(1.25), Inches(4.0), Inches(0.35),
    [[("Shared importance ratio:", 13, GREY, True)]])
pic_fit(s, ASSETS / "eq_ratio.png", Inches(4.4), Inches(1.18), Inches(3.4), Inches(0.7))

# left PPO panel
lx, lw = Inches(0.45), Inches(6.05)
rx = Inches(6.83)
py, ph = Inches(2.1), Inches(4.05)
rect(s, lx, py, lw, ph, PALE_O, line=ORANGE)
rect(s, rx, py, lw, ph, PALE_G, line=GREEN)
panel_title(s, lx, py, lw, "PPO", RUST)
panel_title(s, rx, py, lw, "GRPO", GREEN)

pic_fit(s, ASSETS / "eq_ppo.png", lx + Inches(0.1), py + Inches(0.65), lw - Inches(0.2), Inches(0.95))
pic_fit(s, ASSETS / "eq_ppo_adv.png", lx + Inches(0.1), py + Inches(1.75), lw - Inches(0.2), Inches(0.6))
txt(s, lx + Inches(0.25), py + Inches(2.5), lw - Inches(0.5), Inches(1.4),
    [[("• Advantage Â comes from a learned critic V (GAE).", 14, DARK, False)],
     [("• Extra value network → more memory & compute.", 14, RUST, True)],
     [("• KL penalty keeps the policy near the reference.", 14, DARK, False)]],
    space_after=8, line_spacing=1.05)

pic_fit(s, ASSETS / "eq_grpo.png", rx + Inches(0.1), py + Inches(0.65), lw - Inches(0.2), Inches(0.95))
pic_fit(s, ASSETS / "eq_grpo_adv.png", rx + Inches(0.1), py + Inches(1.75), lw - Inches(0.2), Inches(0.6))
txt(s, rx + Inches(0.25), py + Inches(2.5), lw - Inches(0.5), Inches(1.4),
    [[("• Â = reward normalised within the group (no critic).", 14, DARK, False)],
     [("• Group mean is the baseline → lighter & simpler.", 14, GREEN, True)],
     [("• Same clipped surrogate + KL to the reference.", 14, DARK, False)]],
    space_after=8, line_spacing=1.05)
footer(s, num())

# =============================================================================
# 10 — Results: success rate bars
# =============================================================================
s = slide()
header(s, "Results · Success Rate", "VLM-IN-THE-LOOP vs GROUND-TRUTH BBOX (ORACLE)")
pic_fit(s, ASSETS / "success_bars.png", Inches(0.55), Inches(1.4), Inches(8.45), Inches(5.4))
bullets(s, [
    (0, "Same policy in both runs.", "b"),
    (1, "Only the bbox source changes."),
    (0, "GT bbox peaks at H=2 → 100%.", "ob"),
    (0, "VLM bbox is flatter: 65–78%.", "b"),
    (1, "The gap is the cost of "
        "imperfect perception."),
    (0, "VLM degrades least at H=4–5,", "b"),
    (1, "where open-loop drift hurts "
        "both runs anyway."),
], x=Inches(9.15), y=Inches(1.6), w=Inches(3.9), gap=10, size=14.5)
footer(s, num())

# =============================================================================
# 11 — Results: trends + qualitative
# =============================================================================
s = slide()
header(s, "Results · Trends & A Real Episode", "ACCURACY VS HORIZON, AND WHAT IT LOOKS LIKE")
pic_fit(s, ASSETS / "gt_vs_vlm.png", Inches(0.4), Inches(1.35), Inches(8.4), Inches(3.05))
txt(s, Inches(0.55), Inches(4.45), Inches(8.2), Inches(0.3),
    [[("VLM-bbox episode (H=2): the predicted box (red) tracks the blue cube as the arm reaches it.",
       11.5, GREY, False, True)]])
pic_fit(s, ASSETS / "vlm_episode.png", Inches(0.55), Inches(4.7), Inches(8.1), Inches(2.45))
bullets(s, [
    (0, "Final distance", "ob"),
    (1, "GT reaches ~6 cm at best."),
    (1, "VLM stays ~14–23 cm — bbox "
        "error biases the goal point."),
    (0, "Sweet spot: H = 2–3", "ob"),
    (1, "enough look-ahead, not yet "
        "drifting open-loop."),
    (0, "Takeaway", "ob"),
    (1, "the controller is near-perfect; "
        "perception is the limiter."),
], x=Inches(9.0), y=Inches(1.5), w=Inches(4.05), gap=8, size=14.5)
footer(s, num())

# =============================================================================
# 12 — Results: interpretation
# =============================================================================
s = slide()
header(s, "Results · What the Gap Means", "READING THE NUMBERS")
bullets(s, [
    (0, "The policy is essentially solved.", "ob"),
    (1, "Test RMSE ≈ 0.0007 m, and with oracle bboxes success reaches 100% (H=2). "
        "The MLP is not the bottleneck."),
    (0, "Every VLM-run failure traces back to the bounding box.", "ob"),
    (1, "A box that is off-centre or wrong-sized shifts the goal the arm aims at — "
        "so the arm confidently reaches the wrong place."),
    (1, "Because the bbox is fixed per episode, an early detection error persists "
        "for the whole reach."),
    (0, "Horizon is a perception-cost trade-off.", "ob"),
    (1, "Small H re-queries the VLM often (accurate but slow); large H drifts "
        "open-loop. H=2–3 balances both."),
    (0, "Implication: improving perception — not control — is what raises success. "
        "That is exactly what GRPO fine-tuning would target.", "b"),
], y=Inches(1.5), gap=11, size=18)
footer(s, num())

# =============================================================================
# 13 — Conclusion
# =============================================================================
s = slide()
header(s, "Conclusion & Future Work", "WHAT WE LEARNED")
bullets(s, [
    (0, "A frozen VLM + a tiny behaviour-cloned MLP is enough to build a working "
        "language-conditioned reaching system.", "b"),
    (0, "The learned controller is near-perfect — up to 100% success with oracle "
        "bboxes; it is not the limiting factor.", "ob"),
    (0, "Perception is the bottleneck — swapping the oracle for the off-the-shelf "
        "VLM costs ~15–25 success points.", "ob"),
    (0, "Clear next step:", "ob"),
    (1, "fine-tune the VLM with GRPO (reward = bbox accuracy / reaching success) "
        "to close the perception gap — the one piece compute prevented here."),
    (1, "GRPO is the natural fit: the reward is cheap to compute and it needs no "
        "value network."),
    (0, "Practical operating point: horizon H = 2–3 — strong accuracy with fewer "
        "VLM queries.", "b"),
], y=Inches(1.5), gap=11, size=18)
footer(s, num())

out = HERE / "VLM_Guided_Robot_Reaching.pptx"
prs.save(str(out))
print("saved", out, "slides:", N)
